# ================================
# PPT MCP SERVER (Server 1)
# Runs as stdio subprocess via MCP protocol
# ================================

import os
import sys
import requests
from io import BytesIO
from dotenv import load_dotenv

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from mcp.server.fastmcp import FastMCP

# Load .env relative to this file's location (works as subprocess too)
_here = os.path.dirname(os.path.abspath(__file__))
load_dotenv(os.path.join(_here, '..', '.env'))

mcp = FastMCP("ppt_server")

PEXELS_KEY = os.getenv("PEXELS_API_KEY")

# ================================
# THEME DEFINITIONS
# ================================
THEMES = {
    "Dark Tech": {
        "bg":         RGBColor(0x1a, 0x1a, 0x2e),
        "header":     RGBColor(0xe9, 0x45, 0x60),
        "accent":     RGBColor(0x16, 0x21, 0x3e),
        "title_text": RGBColor(0xff, 0xff, 0xff),
        "body_text":  RGBColor(0xea, 0xea, 0xea),
        "bullet_dot": RGBColor(0xe9, 0x45, 0x60),
    },
    "Corporate Blue": {
        "bg":         RGBColor(0x00, 0x28, 0x55),
        "header":     RGBColor(0x00, 0x88, 0xff),
        "accent":     RGBColor(0x00, 0x44, 0x88),
        "title_text": RGBColor(0xff, 0xff, 0xff),
        "body_text":  RGBColor(0xd0, 0xe8, 0xff),
        "bullet_dot": RGBColor(0x66, 0xcc, 0xff),
    },
    "Nature Green": {
        "bg":         RGBColor(0x14, 0x2b, 0x1e),
        "header":     RGBColor(0x52, 0xb7, 0x88),
        "accent":     RGBColor(0x1e, 0x4d, 0x35),
        "title_text": RGBColor(0xff, 0xff, 0xff),
        "body_text":  RGBColor(0xcc, 0xee, 0xd7),
        "bullet_dot": RGBColor(0x7b, 0xd6, 0xa0),
    },
    "Warm Sunset": {
        "bg":         RGBColor(0x22, 0x11, 0x11),
        "header":     RGBColor(0xff, 0x6b, 0x35),
        "accent":     RGBColor(0x44, 0x22, 0x11),
        "title_text": RGBColor(0xff, 0xff, 0xff),
        "body_text":  RGBColor(0xf5, 0xe0, 0xcc),
        "bullet_dot": RGBColor(0xff, 0xa0, 0x60),
    },
    "Minimal White": {
        "bg":         RGBColor(0xf8, 0xf9, 0xfa),
        "header":     RGBColor(0x4a, 0x90, 0xe2),
        "accent":     RGBColor(0xe8, 0xf0, 0xfe),
        "title_text": RGBColor(0x1a, 0x1a, 0x2e),
        "body_text":  RGBColor(0x33, 0x33, 0x44),
        "bullet_dot": RGBColor(0x4a, 0x90, 0xe2),
    },
}

# Global state (persists across tool calls in same process)
prs = None
current_theme = None
SLIDE_W = Inches(10)
SLIDE_H = Inches(7.5)


# ================================
# INTERNAL HELPERS
# ================================
def _bg_rect(slide, color, left=0, top=0, width=None, height=None):
    w = width or SLIDE_W
    h = height or SLIDE_H
    s = slide.shapes.add_shape(1, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def _textbox(slide, text, left, top, width, height,
             size=18, bold=False, color=None, align=PP_ALIGN.LEFT):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txb


def _place_image(slide, stream, left, top, width, height):
    try:
        if stream:
            stream.seek(0)
            slide.shapes.add_picture(stream, left, top, width=width, height=height)
    except Exception as e:
        print(f"[ppt_server] image error: {e}", file=sys.stderr)


def _fetch_pexels(keyword: str):
    """Fetch image BytesIO from Pexels API."""
    if not PEXELS_KEY or not keyword.strip():
        return None
    try:
        query_param = requests.utils.quote(keyword)
        url = f"https://api.pexels.com/v1/search?query={query_param}&per_page=1&orientation=landscape&size=medium"  # noqa: E501
        r = requests.get(url, headers={"Authorization": PEXELS_KEY}, timeout=10)
        photos = r.json().get("photos", [])
        if photos:
            img_url = photos[0]["src"]["large"]
            img_r = requests.get(img_url, timeout=15)  # noqa: E501
            return BytesIO(img_r.content)
    except Exception as e:
        print(f"[ppt_server] pexels error: {e}", file=sys.stderr)
    return None


# ================================
# INTERNAL SLIDE BUILDERS
# ================================
def _build_title_slide(topic: str, subtitle: str, image_stream):
    t = current_theme
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg_rect(slide, t["bg"])

    if image_stream:
        _bg_rect(slide, t["accent"],
                 left=Inches(5.8), top=Inches(0),
                 width=Inches(4.2), height=SLIDE_H)
        _place_image(slide, image_stream,
                     Inches(5.9), Inches(0.3), Inches(4.0), Inches(6.9))

    _bg_rect(slide, t["header"],
             left=0, top=0, width=Inches(0.18), height=SLIDE_H)
    _bg_rect(slide, t["header"],
             left=0, top=Inches(7.15), width=SLIDE_W, height=Inches(0.12))
    _bg_rect(slide, t["accent"],
             left=Inches(0.3), top=Inches(4.0),
             width=Inches(5.2), height=Inches(0.06))

    _textbox(slide, topic,
             Inches(0.4), Inches(1.6), Inches(5.4), Inches(1.8),
             size=38, bold=True, color=t["title_text"], align=PP_ALIGN.LEFT)
    _textbox(slide, subtitle,
             Inches(0.4), Inches(4.2), Inches(5.2), Inches(0.7),
             size=16, bold=False, color=t["bullet_dot"], align=PP_ALIGN.LEFT)


def _build_content_slide(title: str, bullets: list, image_stream):
    t = current_theme
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg_rect(slide, t["bg"])

    has_img = image_stream is not None
    txt_width = Inches(5.5) if has_img else Inches(9.5)

    if has_img:
        _bg_rect(slide, t["accent"],
                 left=Inches(6.1), top=Inches(1.25),
                 width=Inches(3.6), height=Inches(6.0))
        _place_image(slide, image_stream,
                     Inches(6.15), Inches(1.3), Inches(3.5), Inches(5.9))

    _bg_rect(slide, t["header"],
             left=0, top=0, width=SLIDE_W, height=Inches(1.2))
    _bg_rect(slide, t["bullet_dot"],
             left=0, top=Inches(1.2), width=SLIDE_W, height=Inches(0.05))

    _textbox(slide, title,
             Inches(0.25), Inches(0.1), Inches(9.3), Inches(1.0),
             size=26, bold=True, color=RGBColor(0xff, 0xff, 0xff),
             align=PP_ALIGN.LEFT)

    if bullets:
        txb = slide.shapes.add_textbox(
            Inches(0.3), Inches(1.45), txt_width, Inches(5.8))
        tf = txb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            if i > 0:
                p.space_before = Pt(24)
            dot = p.add_run()
            dot.text = "● "
            dot.font.size = Pt(14)
            dot.font.color.rgb = t["bullet_dot"]
            run = p.add_run()
            run.text = bullet.strip("- •▶●").strip()
            run.font.size = Pt(22)
            run.font.color.rgb = t["body_text"]

    _bg_rect(slide, t["header"],
             left=0, top=Inches(7.38), width=SLIDE_W, height=Inches(0.12))


def _build_conclusion_slide(topic: str):
    t = current_theme
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg_rect(slide, t["bg"])
    _bg_rect(slide, t["accent"],
             left=Inches(1.5), top=Inches(1.8),
             width=Inches(7.0), height=Inches(4.0))
    _bg_rect(slide, t["header"],
             left=Inches(1.5), top=Inches(1.8),
             width=Inches(7.0), height=Inches(0.18))
    _textbox(slide, "Thank You! 🙏",
             Inches(1.5), Inches(2.3), Inches(7.0), Inches(1.4),
             size=42, bold=True, color=t["title_text"], align=PP_ALIGN.CENTER)
    _textbox(slide, f"Topic: {topic}",
             Inches(1.5), Inches(3.8), Inches(7.0), Inches(0.7),
             size=17, color=t["bullet_dot"], align=PP_ALIGN.CENTER)
    _bg_rect(slide, t["header"],
             left=0, top=Inches(7.15), width=SLIDE_W, height=Inches(0.12))


# ================================
# MCP TOOL 1: create_presentation
# ================================
@mcp.tool()
def create_presentation(theme_name: str = "Dark Tech") -> str:
    """
    Initialize a new PowerPoint presentation with the given theme.

    Args:
        theme_name: One of 'Dark Tech', 'Corporate Blue', 'Nature Green',  # noqa: E501
                    'Warm Sunset', 'Minimal White'
    Returns:
        Confirmation string
    """
    global prs, current_theme  # noqa: F824
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    current_theme = THEMES.get(theme_name, THEMES["Dark Tech"])
    return f"✅ Presentation initialized with theme: {theme_name}"


# ================================
# MCP TOOL 2: add_slide
# ================================
@mcp.tool()
def add_slide(title: str, content: str,
              image_keyword: str = "", slide_type: str = "content") -> str:
    """
    Add a single slide to the presentation.

    Args:
        title: Slide title text
        content: Bullet points (content slides) or subtitle (title slides)
        image_keyword: Pexels search term for slide image (empty = no image)  # noqa: E501
        slide_type: 'title' | 'content' | 'conclusion'  # noqa: E501
    Returns:
        Confirmation string
    """
    global prs, current_theme  # noqa: F824
    if prs is None:
        return "❌ Error: call create_presentation() first"

    # Fetch image via Pexels (inside the server, keeps MCP params JSON-serializable)
    image_stream = _fetch_pexels(image_keyword) if image_keyword.strip() else None

    if slide_type == "title":
        _build_title_slide(title, content, image_stream)
        return f"✅ Title slide added: '{title}'"

    elif slide_type == "conclusion":
        _build_conclusion_slide(title)
        return "✅ Conclusion slide added"

    else:  # content slide
        bullets = [b.strip() for b in content.split("\n") if b.strip()]
        _build_content_slide(title, bullets, image_stream)
        return f"✅ Content slide added: '{title}' ({len(bullets)} bullets)"


# ================================
# MCP TOOL 3: save_presentation
# ================================
@mcp.tool()
def save_presentation(file_path: str) -> str:
    """
    Save the presentation to disk.

    Args:
        file_path: Output path for the .pptx file
    Returns:
        Confirmation string
    """
    global prs
    if prs is None:
        return "❌ Error: no presentation to save"
    try:
        os.makedirs(os.path.dirname(os.path.abspath(file_path)), exist_ok=True)
        if os.path.exists(file_path):
            os.remove(file_path)
        prs.save(file_path)
        return f"✅ Presentation saved: {file_path}"
    except Exception as e:
        return f"❌ Save error: {str(e)}"


# ================================
# RUN AS MCP STDIO SERVER
# ================================
if __name__ == "__main__":
    mcp.run()
